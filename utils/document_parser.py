"""
Document parsing utilities.
Extracts text + page/slide-level metadata from PDF, DOCX, and PPTX files.

Each function returns a list of dicts:
    {
        "text": str,           # raw extracted text for that unit (page/slide/paragraph-block)
        "source": str,         # original filename
        "page": int | None,    # page number (PDF) or slide number (PPTX)
        "section": str | None  # heading/section label if available (DOCX)
    }
"""

import io
import pdfplumber
import docx
from pptx import Presentation


def parse_pdf(file_bytes: bytes, filename: str) -> list[dict]:
    """Extract text from a PDF, one record per page."""
    records = []
    with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
        for i, page in enumerate(pdf.pages, start=1):
            text = page.extract_text() or ""
            text = text.strip()
            if text:
                records.append({
                    "text": text,
                    "source": filename,
                    "page": i,
                    "section": None,
                })
    return records


def parse_docx(file_bytes: bytes, filename: str) -> list[dict]:
    """
    Extract text from a DOCX file.
    DOCX has no native page numbers (pagination is a rendering concept),
    so we track the most recent Heading as the 'section' and chunk by
    paragraph groups instead of pages.
    """
    document = docx.Document(io.BytesIO(file_bytes))
    records = []
    current_section = "Document Start"
    buffer = []

    def flush():
        text = "\n".join(buffer).strip()
        if text:
            records.append({
                "text": text,
                "source": filename,
                "page": None,
                "section": current_section,
            })
        buffer.clear()

    for para in document.paragraphs:
        style_name = (para.style.name or "").lower()
        text = para.text.strip()
        if not text:
            continue
        if style_name.startswith("heading"):
            # New section begins -> flush what we have under the old section
            flush()
            current_section = text
        else:
            buffer.append(text)
            # Flush every ~6 paragraphs to keep chunks reasonably sized
            if len(buffer) >= 6:
                flush()
    flush()
    return records


def parse_pptx(file_bytes: bytes, filename: str) -> list[dict]:
    """Extract text from a PPTX file, one record per slide."""
    prs = Presentation(io.BytesIO(file_bytes))
    records = []
    for i, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs).strip()
                    if line:
                        texts.append(line)
            # Also pull text from tables on the slide
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            texts.append(cell.text.strip())
        text = "\n".join(texts).strip()
        if text:
            records.append({
                "text": text,
                "source": filename,
                "page": i,  # "page" doubles as slide number here
                "section": None,
            })
    return records


def parse_file(file_bytes: bytes, filename: str) -> list[dict]:
    """Dispatch to the correct parser based on file extension."""
    ext = filename.lower().rsplit(".", 1)[-1]
    if ext == "pdf":
        return parse_pdf(file_bytes, filename)
    elif ext == "docx":
        return parse_docx(file_bytes, filename)
    elif ext in ("pptx", "ppt"):
        return parse_pptx(file_bytes, filename)
    else:
        raise ValueError(f"Unsupported file type: {ext}")
